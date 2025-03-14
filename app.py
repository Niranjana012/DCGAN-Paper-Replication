from flask import Flask, render_template, request, send_file
import torch
import torchvision.utils as vutils
import os
from pptx import Presentation
from pptx.util import Inches
from generator import Generator

app = Flask(__name__)

# Load Generator Model
latent_dim = 100
device = torch.device("cuda" if torch.cuda.is_available() else "cpu")

MODEL_PATH = "generator_epoch_100.pth"
if not os.path.exists(MODEL_PATH):
    print(f"❌ Model file {MODEL_PATH} not found. Train the model first!")
    exit()

generator = Generator(latent_dim).to(device)
generator.load_state_dict(torch.load(MODEL_PATH, map_location=device))
generator.eval()

# Route for home page
@app.route('/')
def home():
    return render_template('index.html')

# Route for generating images
@app.route('/generate', methods=['POST'])
def generate():
    num_images = int(request.form.get('num_images', 10))
    os.makedirs("static/generated_samples", exist_ok=True)
    file_paths = []
    
    with torch.no_grad():
        for i in range(num_images):
            noise = torch.randn(1, latent_dim, 1, 1, device=device)
            fake_image = generator(noise).cpu()
            img_path = f"static/generated_samples/generated_slide_{i}.png"
            vutils.save_image(fake_image, img_path, normalize=True)
            file_paths.append(img_path)
    
    return render_template('result.html', images=file_paths)

# Route for downloading all images as a PPTX
@app.route('/download_pptx')
def download_pptx():
    ppt = Presentation()
    image_folder = "static/generated_samples"
    image_files = sorted([f for f in os.listdir(image_folder) if f.endswith(".png")])

    if not image_files:
        return "No generated images found."

    for img_file in image_files:
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Blank slide
        img_path = os.path.join(image_folder, img_file)
        left = top = Inches(1)
        slide.shapes.add_picture(img_path, left, top, width=Inches(8))

    ppt_output = "static/generated_slides.pptx"
    ppt.save(ppt_output)
    return send_file(ppt_output, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
